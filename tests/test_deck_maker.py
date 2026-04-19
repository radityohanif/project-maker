from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation

from deck_maker.core.images import resolve_image_path
from deck_maker.core.models import DeckImageSource, ImageSlide, PresentationSpec
from deck_maker.core.prompt_builder import PromptParams, build_ai_prompt
from deck_maker.core.parser import parse_file
from deck_maker.core.renderer import render


def test_remote_image_requires_allow_network(tmp_path: Path) -> None:
    spec = PresentationSpec(
        allow_network=False,
        slides=[
            ImageSlide(
                title="Fig",
                source=DeckImageSource(url="https://placehold.co/10x10"),
            )
        ],
    )
    out = tmp_path / "a.pptx"
    with pytest.raises(ValueError, match="allow_network"):
        render(spec, out, base_dir=tmp_path)


def test_remote_image_allowed_with_mock_fetch(tmp_path: Path) -> None:
    """URL fetch succeeds when allow_network is true (HTTP mocked)."""
    png = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00"
        b"\x00\x01\x01\x00\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    mock_resp = MagicMock()
    mock_resp.read.return_value = png
    mock_cm = MagicMock()
    mock_cm.__enter__.return_value = mock_resp
    mock_cm.__exit__.return_value = None

    spec = PresentationSpec(
        allow_network=True,
        slides=[
            ImageSlide(
                title="Fig",
                source=DeckImageSource(url="https://example.test/x.png"),
            )
        ],
    )
    out = tmp_path / "b.pptx"
    with patch("deck_maker.core.images.urllib.request.urlopen", return_value=mock_cm):
        render(spec, out, base_dir=tmp_path)
    assert out.is_file()
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_parse_and_render_standalone_deck(tmp_path: Path) -> None:
    deck_yaml = tmp_path / "deck.yaml"
    deck_yaml.write_text(
        """
title: "Standalone"
slides:
  - type: title
    title: "Hello"
    subtitle: "Deck"
""",
        encoding="utf-8",
    )
    spec = parse_file(deck_yaml)
    out = tmp_path / "out.pptx"
    render(spec, out, base_dir=tmp_path)
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_resolve_image_path_local(tmp_path: Path) -> None:
    img = tmp_path / "x.png"
    img.write_bytes(b"not a real png but ok for path test")
    src = DeckImageSource(path="x.png")
    got = resolve_image_path(src, base_dir=tmp_path, allow_network=False)
    assert got == img.resolve()


def test_resolve_image_path_url_blocked() -> None:
    src = DeckImageSource(url="https://example.com/x.png")
    with pytest.raises(ValueError, match="allow_network"):
        resolve_image_path(src, base_dir=Path("."), allow_network=False)


def test_build_ai_prompt_includes_deck_schema() -> None:
    text = build_ai_prompt(
        PromptParams(
            deck_title="Pitch",
            topic_hint="Series A",
            allow_network=False,
            output_name="deck.pptx",
            template_path="",
            slides_hint="title; bullets",
            strict_markdown=True,
        )
    )
    assert "deck-maker" in text
    assert "Pitch" in text
    assert "allow_network" in text
    assert "type: table" in text or "`type: table`" in text
