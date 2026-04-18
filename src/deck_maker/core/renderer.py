from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from deck_maker.core.images import cleanup_temp_remote, resolve_image_path
from deck_maker.core.models import (
    BulletsSlide,
    ImageSlide,
    PresentationSpec,
    SectionSlide,
    TableSlide,
    TitleSlide,
)


def _find_layout(prs: Presentation, *candidates: str) -> object:
    lowered = [c.lower() for c in candidates]
    for layout in prs.slide_layouts:
        name = layout.name.lower()
        if name in lowered:
            return layout
        for cand in lowered:
            if cand in name:
                return layout
    return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]


def _cell_to_str(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def render(
    spec: PresentationSpec,
    path: Path,
    *,
    base_dir: Path,
) -> Path:
    """Write a ``.pptx`` from ``spec``. ``base_dir`` resolves relative image paths."""
    if spec.template:
        tpl = (base_dir / spec.template).resolve()
        if not tpl.is_file():
            raise FileNotFoundError(f"Deck template not found: {tpl}")
        prs = Presentation(str(tpl))
    else:
        prs = Presentation()

    layout_title = _find_layout(prs, "title slide", "title")
    layout_content = _find_layout(prs, "title and content", "title and content")
    layout_section = _find_layout(prs, "section header", "section")
    layout_blank = _find_layout(prs, "blank")

    for block in spec.slides:
        if isinstance(block, TitleSlide):
            slide = prs.slides.add_slide(layout_title)
            if slide.shapes.title:
                slide.shapes.title.text = block.title
            sub_parts: list[str] = []
            if block.subtitle:
                sub_parts.append(block.subtitle)
            if block.meta:
                sub_parts.append(block.meta)
            if sub_parts and len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                ph.text = "\n".join(sub_parts)
        elif isinstance(block, SectionSlide):
            slide = prs.slides.add_slide(layout_section)
            if slide.shapes.title:
                slide.shapes.title.text = block.title
        elif isinstance(block, BulletsSlide):
            slide = prs.slides.add_slide(layout_content)
            if slide.shapes.title:
                slide.shapes.title.text = block.title
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            if body is not None:
                tf = body.text_frame
                tf.clear()
                items = block.bullets or [""]
                tf.text = items[0]
                p0 = tf.paragraphs[0]
                p0.font.size = Pt(18)
                for item in items[1:]:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
        elif isinstance(block, TableSlide):
            slide = prs.slides.add_slide(layout_content)
            if slide.shapes.title:
                slide.shapes.title.text = block.title or ""
            rows = 1 + len(block.rows)
            cols = len(block.headers)
            if cols == 0:
                continue
            left, top, width = Inches(0.5), Inches(1.5), Inches(9.0)
            row_h = Inches(0.35)
            height = row_h * rows
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            tbl = table_shape.table
            for c, h in enumerate(block.headers):
                tbl.cell(0, c).text = str(h)
            for r, row in enumerate(block.rows, start=1):
                for c in range(cols):
                    val = row[c] if c < len(row) else ""
                    tbl.cell(r, c).text = _cell_to_str(val)
        elif isinstance(block, ImageSlide):
            slide = prs.slides.add_slide(layout_blank)
            y = Inches(0.35)
            if block.title:
                title_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9.0), Inches(0.7))
                title_box.text_frame.text = block.title
                y = Inches(1.15)
            img_path = resolve_image_path(
                block.source, base_dir=base_dir, allow_network=spec.allow_network
            )
            remote_tmp = block.source.url is not None
            try:
                slide.shapes.add_picture(str(img_path), Inches(1), y, width=Inches(7.5))
            finally:
                if remote_tmp:
                    cleanup_temp_remote(img_path)
            if block.caption:
                cap_top = Inches(6.5)
                cap = slide.shapes.add_textbox(Inches(0.75), cap_top, Inches(8.5), Inches(0.6))
                cap.text_frame.text = block.caption
                for para in cap.text_frame.paragraphs:
                    para.font.size = Pt(11)
                    para.font.italic = True

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
